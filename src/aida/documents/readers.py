"""Structured document reading, dispatched by file extension (PLAN.md Phase
6): "Dispatcher by extension/MIME; structured extraction, not blind
flattening." Each format gets format-appropriate extraction into AIDA's
typed ``Artifact`` types (PLAN.md hard rule 3) rather than one giant
"read bytes, hope it's text" fallback:

- Plain text / code / Markdown -> ``TextArtifact``
- CSV -> ``TableArtifact`` (stdlib ``csv``)
- JSON -> ``JsonArtifact`` (structure preserved, not re-flattened to text)
- PDF (extra ``docs``, ``pymupdf``) -> ``TextArtifact``, one page per
  ``--- page N ---`` section, capped at ``max_pdf_pages``
- DOCX (extra ``docs``, ``python-docx``) -> ``TextArtifact`` of paragraph
  text (tables rendered inline as pipe rows)
- XLSX (extra ``docs``, ``openpyxl``) -> one ``TableArtifact`` per sheet,
  capped at ``max_sheets`` sheets / ``max_rows_per_sheet`` rows each
- PPTX (extra ``docs``, ``python-pptx``) -> ``TextArtifact``, one slide per
  ``--- slide N ---`` section
- Images -> ``ImageArtifact`` referencing the file's own path (no bytes
  loaded into memory — the GUI's ``InlineImageWidget`` reads straight off
  disk via the path, same as a live ``ImageArtifactCreated`` event)

**Vision input (B1).** An ``ImageArtifact`` read this way (a GUI drag-and-
drop or "Attach…" image) is described in text (via ``describe_for_model``,
as every artifact type is) *and*, when the active profile has
``supports_vision: true``, has its actual pixels attached as vision input —
see ``aida.ui.qt.main_window._augment_with_attachments`` and
``aida.providers.vision``. ``aida.providers.base.Message.content`` itself
stays plain ``str`` throughout the provider layer; the image travels
alongside it on ``Message.images`` rather than being folded into content.
Without ``supports_vision`` set (the default — not every endpoint AIDA
talks to understands image content blocks), the image still displays in
the GUI and is still described in text, exactly as before B1.

HDF5 is deliberately **not** implemented here — that's pyIrena MCP's job
(Phase 3 already covers science-data formats via MCP), not a general
document reader.

**None of these readers extract images** — every one is text-only, and a
figure, plot, scanned page or embedded chart is dropped. Since a silently
incomplete document is worse than an obviously incomplete one, each reader
appends a short note saying how many images it dropped, and the PDF reader
additionally detects the case where there was no usable text layer at all
(a scanned paper, a photographed logbook) and says so rather than handing
back an empty-looking document. See ``planning/document_images.md`` for the
richer options this note is the floor of.

Every reader applies a size/token guard (PLAN.md: "long docs summarized-by-
section or chunk-selected rather than context-bombed") — callers that want
the *whole* document regardless of size should read the file directly
rather than going through this dispatcher.
"""

from __future__ import annotations

import csv
import json
import mimetypes
import zipfile
from collections.abc import Callable
from pathlib import Path

from aida.artifacts.base import Artifact, ImageArtifact, JsonArtifact, TableArtifact, TextArtifact

DEFAULT_MAX_CHARS = 20_000
DEFAULT_MAX_PDF_PAGES = 50
DEFAULT_MAX_SHEETS = 5
DEFAULT_MAX_ROWS_PER_SHEET = 200
DEFAULT_MAX_CSV_ROWS = 500

#: Budget for the two *interactive* read paths — the `read_file` tool and
#: the GUI drag-and-drop/"Attach..." path — where the point is handing a
#: real document (e.g. a journal paper) to the model in one shot, as
#: opposed to `DEFAULT_MAX_CHARS` above (sized for cheap fallback text
#: reads) or RAG ingestion's much larger `_INGEST_MAX_CHARS` (which wants
#: the whole document for chunking, not chat context). Callers using these
#: must also pass a matching `max_chars` to `describe_for_model()` —
#: otherwise its own smaller default silently re-truncates on top of this.
INTERACTIVE_MAX_CHARS = 100_000
INTERACTIVE_MAX_PDF_PAGES = 150

_TRUNCATION_NOTE = "\n... [truncated]"

_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}
_TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".markdown",
    ".rst",
    ".log",
    ".py",
    ".js",
    ".ts",
    ".jsx",
    ".tsx",
    ".java",
    ".c",
    ".h",
    ".cpp",
    ".hpp",
    ".rs",
    ".go",
    ".rb",
    ".sh",
    ".bash",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
    ".cfg",
    ".xml",
    ".html",
    ".htm",
    ".css",
    ".sql",
}

#: Below this many non-whitespace characters *per page*, a PDF is treated as
#: having no usable text layer. The failure this catches: a scanned paper, a
#: photographed logbook page or a signed form extracts as empty or
#: near-empty text, and the model receives what looks like an empty document
#: with no way to tell that from an unreadable one — so it guesses, or
#: reports the attachment as blank. 50 is deliberately generous: a real page
#: of prose is thousands of characters, while a scanned page still yields a
#: few stray glyphs from headers, stamps or OCR already baked into the file.
_SCANNED_TEXT_CHARS_PER_PAGE = 50

#: Where each Office format keeps its embedded media inside its own zip
#: container. Counting entries here is stdlib-only, cheap (the central
#: directory, not the file bodies), independent of library internals, and —
#: unlike ``openpyxl``'s ``sheet._images`` — works with the ``read_only``
#: loading mode the xlsx reader deliberately uses.
_MEDIA_DIR_BY_SUFFIX = {
    ".docx": "word/media/",
    ".xlsx": "xl/media/",
    ".pptx": "ppt/media/",
}


def _count_embedded_media(path: Path) -> int:
    """How many embedded media files an Office document contains.

    Slightly over-counts for ``.pptx``, where ``ppt/media/`` also holds
    images belonging to slide layouts, masters and the theme. That is the
    right direction to be wrong in: the note this feeds tells the model
    content was dropped, and over-reporting makes it ask, while
    under-reporting would let it assume a figure-heavy deck was fully read.
    """
    prefix = _MEDIA_DIR_BY_SUFFIX.get(path.suffix.lower())
    if prefix is None:
        return 0
    try:
        with zipfile.ZipFile(path) as archive:
            return sum(
                1
                for name in archive.namelist()
                if name.startswith(prefix) and not name.endswith("/")
            )
    except (OSError, zipfile.BadZipFile):
        # Not fatal: the reader itself has already parsed this file, so a
        # failure here means only that the count is unavailable. Say
        # nothing rather than guess.
        return 0


def _dropped_images_note(count: int) -> str:
    """The note appended when a reader extracted text but silently dropped
    the document's images — every reader in this module is text-only (see
    the module docstring). Empty string when there is nothing to say."""
    if count <= 0:
        return ""
    subject = (
        "1 embedded image, which was" if count == 1 else f"{count} embedded images, which were"
    )
    return (
        f"\n\n[This document contains {subject} not extracted — only its text "
        f"is shown above. If it was attached to this conversation, "
        f"list_document_figures can name its figures and get_document_figure "
        f"can show you one; otherwise say the answer depends on a figure you "
        f"cannot see.]"
    )


def _pdf_content_note(*, pages_read: int, text_chars: int, image_count: int) -> str:
    """The PDF equivalent, which also has to cover the case a plain image
    count cannot: a page-image-only document, where the *absence* of text is
    the finding rather than the presence of figures."""
    if pages_read and text_chars < _SCANNED_TEXT_CHARS_PER_PAGE * pages_read:
        if image_count:
            pages = "1 page" if pages_read == 1 else f"{pages_read} pages"
            images = (
                "1 image, which was" if image_count == 1 else f"{image_count} images, which were"
            )
            return (
                f"\n\n[No usable text layer — this appears to be a scanned or "
                f"image-only PDF. Its {pages} hold {images} not extracted. "
                f"Treat this as a document that could not be read, not as an "
                f"empty one.]"
            )
        return (
            "\n\n[No extractable text and no embedded images — this PDF may be "
            "empty or damaged. Treat it as unread rather than blank.]"
        )
    return _dropped_images_note(image_count)


class UnsupportedDocumentFormatError(Exception):
    """Raised for a file extension no reader (and no plain-text fallback)
    recognizes — surfaces as a clear tool error rather than either crashing
    or silently mis-decoding binary data as text."""


def _truncate(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[: max_chars - len(_TRUNCATION_NOTE)] + _TRUNCATION_NOTE


def _read_text_file(path: Path, *, max_chars: int) -> list[Artifact]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [TextArtifact(text=_truncate(text, max_chars))]


def _read_csv_file(
    path: Path, *, max_rows: int = DEFAULT_MAX_CSV_ROWS, **_ignored
) -> list[Artifact]:
    with path.open("r", encoding="utf-8", errors="replace", newline="") as fh:
        reader = csv.reader(fh)
        rows = list(reader)
    if not rows:
        return [TableArtifact(columns=[], rows=[])]
    columns, data_rows = rows[0], rows[1:]
    truncated = data_rows[:max_rows]
    if len(data_rows) > max_rows:
        truncated.append(
            [f"... [{len(data_rows) - max_rows} more rows truncated]"] + [""] * (len(columns) - 1)
        )
    return [TableArtifact(columns=columns, rows=truncated)]


def _read_json_file(path: Path, *, max_chars: int, **_ignored) -> list[Artifact]:
    text = path.read_text(encoding="utf-8", errors="replace")
    try:
        data = json.loads(text)
    except ValueError:
        # Not actually valid JSON despite the extension -- fall back to text
        # rather than raising, since the content still exists and is useful.
        return [TextArtifact(text=_truncate(text, max_chars))]
    return [JsonArtifact(data=data)]


def _read_image_file(path: Path, **_ignored) -> list[Artifact]:
    mime_type = mimetypes.guess_type(str(path))[0] or "application/octet-stream"
    return [ImageArtifact(data=b"", mime_type=mime_type, path=str(path), filename=path.name)]


def _read_pdf_file(
    path: Path, *, max_chars: int, max_pdf_pages: int = DEFAULT_MAX_PDF_PAGES, **_ignored
) -> list[Artifact]:
    import pymupdf

    parts: list[str] = []
    image_xrefs: set[int] = set()
    text_chars = 0
    pages_read = 0
    with pymupdf.open(path) as doc:
        page_count = doc.page_count
        for index, page in enumerate(doc):
            if index >= max_pdf_pages:
                parts.append(f"--- [{page_count - max_pdf_pages} more pages truncated] ---")
                break
            page_text = page.get_text()
            parts.append(f"--- page {index + 1} ---\n{page_text}")
            pages_read += 1
            # Non-whitespace only, and counted from the page text rather
            # than from `parts` — the "--- page N ---" markers are ours, and
            # counting them would make an image-only PDF look like it had
            # text in proportion to its page count.
            text_chars += len("".join(page_text.split()))
            # Keyed by xref so one logo repeated on every page counts once,
            # not fourteen times. Metadata only; no pixels are decoded.
            image_xrefs.update(image[0] for image in page.get_images(full=False))
    body = _truncate("\n\n".join(parts), max_chars)
    # Appended *after* truncation, deliberately: a note explaining what was
    # dropped is worthless if it is itself the thing that gets dropped. It
    # costs a few dozen characters over the budget in the worst case.
    return [
        TextArtifact(
            text=body
            + _pdf_content_note(
                pages_read=pages_read, text_chars=text_chars, image_count=len(image_xrefs)
            )
        )
    ]


def _read_docx_file(path: Path, *, max_chars: int, **_ignored) -> list[Artifact]:
    import docx

    document = docx.Document(str(path))
    parts: list[str] = [p.text for p in document.paragraphs]
    for table in document.tables:
        parts.append("")
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    body = _truncate("\n".join(parts), max_chars)
    return [TextArtifact(text=body + _dropped_images_note(_count_embedded_media(path)))]


def _read_xlsx_file(
    path: Path,
    *,
    max_sheets: int = DEFAULT_MAX_SHEETS,
    max_rows_per_sheet: int = DEFAULT_MAX_ROWS_PER_SHEET,
    **_ignored,
) -> list[Artifact]:
    import openpyxl

    workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        artifacts: list[Artifact] = []
        for sheet_name in workbook.sheetnames[:max_sheets]:
            sheet = workbook[sheet_name]
            all_rows = list(sheet.iter_rows(values_only=True))
            if not all_rows:
                artifacts.append(TableArtifact(columns=[sheet_name], rows=[]))
                continue
            header = [str(c) if c is not None else "" for c in all_rows[0]]
            data_rows = [
                [("" if c is None else c) for c in row]
                for row in all_rows[1 : max_rows_per_sheet + 1]
            ]
            if len(all_rows) - 1 > max_rows_per_sheet:
                data_rows.append(
                    [f"... [{len(all_rows) - 1 - max_rows_per_sheet} more rows truncated]"]
                    + [""] * (len(header) - 1)
                )
            artifacts.append(
                TableArtifact(
                    columns=[f"{sheet_name}: {c}" if header else sheet_name for c in header]
                    or [sheet_name],
                    rows=data_rows,
                )
            )
        if len(workbook.sheetnames) > max_sheets:
            artifacts.append(
                TextArtifact(
                    text=f"... [{len(workbook.sheetnames) - max_sheets} more sheet(s) truncated]"
                )
            )
        # Its own artifact rather than appended to a table's text: these are
        # TableArtifacts, whose whole point is that they are not flattened
        # into prose. The truncation marker above already sets the pattern.
        note = _dropped_images_note(_count_embedded_media(path))
        if note:
            artifacts.append(TextArtifact(text=note.strip()))
        return artifacts
    finally:
        workbook.close()


def _read_pptx_file(path: Path, *, max_chars: int, **_ignored) -> list[Artifact]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides):
        texts = [
            shape.text_frame.text
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text
        ]
        parts.append(f"--- slide {index + 1} ---\n" + "\n".join(texts))
    body = _truncate("\n\n".join(parts), max_chars)
    return [TextArtifact(text=body + _dropped_images_note(_count_embedded_media(path)))]


_READERS: dict[str, Callable[..., list[Artifact]]] = {
    ".csv": _read_csv_file,
    ".json": _read_json_file,
    ".pdf": _read_pdf_file,
    ".docx": _read_docx_file,
    ".xlsx": _read_xlsx_file,
    ".pptx": _read_pptx_file,
}


def is_supported(path: str | Path) -> bool:
    suffix = Path(path).suffix.lower()
    return (
        suffix in _READERS or suffix in _TEXT_SUFFIXES or suffix in _IMAGE_SUFFIXES or suffix == ""
    )


def is_image_path(path: str | Path) -> bool:
    """Whether ``path`` is one of the image formats this module reads as an
    ``ImageArtifact`` (B1: the single source of truth for "is this an
    image attachment" — used by ``aida.ui.qt.main_window`` to decide
    whether a GUI-attached file should also be sent as vision input, not
    just described in text)."""
    return Path(path).suffix.lower() in _IMAGE_SUFFIXES


def read_document(
    path: str | Path,
    *,
    max_chars: int = DEFAULT_MAX_CHARS,
    max_pdf_pages: int = DEFAULT_MAX_PDF_PAGES,
    max_sheets: int = DEFAULT_MAX_SHEETS,
    max_rows_per_sheet: int = DEFAULT_MAX_ROWS_PER_SHEET,
) -> list[Artifact]:
    """Reads ``path`` into one or more typed ``Artifact``s, dispatched by
    file extension. Returns a *list* deliberately (not one ``Artifact``):
    a multi-sheet spreadsheet reads as one ``TableArtifact`` per sheet.

    Raises ``UnsupportedDocumentFormatError`` for an extension with no
    reader and no plain-text fallback (binary formats like ``.zip``,
    ``.exe``, unrecognized proprietary formats, ...) rather than guessing.
    """
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix in _IMAGE_SUFFIXES:
        return _read_image_file(p)
    if suffix in _READERS:
        return _READERS[suffix](
            p,
            max_chars=max_chars,
            max_pdf_pages=max_pdf_pages,
            max_sheets=max_sheets,
            max_rows_per_sheet=max_rows_per_sheet,
        )
    if suffix in _TEXT_SUFFIXES or suffix == "":
        # No extension (README, Makefile, ...) is treated as plain text
        # too, same as any shell would.
        return _read_text_file(p, max_chars=max_chars)

    raise UnsupportedDocumentFormatError(
        f"No reader for {suffix!r} files ({p.name}). Supported: text/code/Markdown, "
        ".csv, .json, .pdf, .docx, .xlsx, .pptx, and common image formats."
    )


__all__ = [
    "DEFAULT_MAX_CHARS",
    "DEFAULT_MAX_CSV_ROWS",
    "DEFAULT_MAX_PDF_PAGES",
    "DEFAULT_MAX_ROWS_PER_SHEET",
    "DEFAULT_MAX_SHEETS",
    "INTERACTIVE_MAX_CHARS",
    "INTERACTIVE_MAX_PDF_PAGES",
    "UnsupportedDocumentFormatError",
    "is_image_path",
    "is_supported",
    "read_document",
]
